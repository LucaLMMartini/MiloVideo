from __future__ import annotations

import json
import logging
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .config import settings
from .schemas import FactSheet, Job

log = logging.getLogger(__name__)

GRAY = RGBColor(0x80, 0x80, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0x0F, 0x62, 0xFE)

# Theme buckets used when no LLM is available (keyword → section title).
_FALLBACK_THEMES = [
    ("Lenkrad & Bedienung", ["lenkrad", "steering", "wheel", "paddle", "schalt", "hupe", "blinker", "wiper", "scheibenwischer"]),
    ("Konnektivität & Anschlüsse", ["usb", "bluetooth", "wlan", "wifi", "carplay", "android auto", "anschluss", "port", "laden", "charging", "wireless", "sim", "hotspot", "nfc", "digital key"]),
    ("Display & Infotainment / UI", ["display", "screen", "bildschirm", "touch", "menü", "menu", "app", "icon", "ui", "infotainment", "widget", "startseite", "home", "navigation", "karte", "map"]),
    ("Klima & Komfort", ["klima", "climate", "temperatur", "heiz", "lüft", "vent", "ac", "defrost", "sitzheizung", "ambient", "ambiente", "beleuchtung", "light"]),
    ("Sitze & Innenraum", ["sitz", "seat", "armlehne", "cupholder", "becher", "handschuh", "glovebox", "ablage", "konsole", "console"]),
    ("Fahrassistenz & Fahrmodi", ["assist", "spur", "lane", "tempomat", "cruise", "park", "kamera", "camera", "sensor", "modus", "mode", "sport", "comfort", "eco", "auto hold", "traction"]),
    ("Instrumente & Anzeigen", ["tacho", "cluster", "instrument", "kombiinstrument", "anzeige", "gauge", "head-up", "hud"]),
]


def _fmt_ts(t: float) -> str:
    m, s = divmod(int(t), 60)
    h, m = divmod(m, 60)
    return f"{h:02d}:{m:02d}:{s:02d}" if h else f"{m:02d}:{s:02d}"


def _items_from_sheet(sheet: FactSheet) -> list[dict]:
    """Flatten facts + features into an indexed list, skipping rejected ones."""
    items: list[dict] = []
    for f in sheet.atomic_facts:
        if f.status == "rejected":
            continue
        items.append({
            "kind": "fact",
            "text": f.fact,
            "evidence": [e.model_dump() for e in f.evidence],
        })
    for f in sheet.features:
        if f.status == "rejected":
            continue
        label = f.label + (f": {f.description}" if f.description else "")
        items.append({
            "kind": "feature",
            "text": label,
            "evidence": [e.model_dump() for e in f.evidence],
        })
    return items


def _vehicle_label(job: Job, sheet: FactSheet) -> str:
    """Prefer the user-provided brand/model/trim, fall back to the detected model."""
    user = " ".join(p for p in (job.brand, job.model_name, job.trim) if p).strip()
    return user or (sheet.vehicle_model or "")


def _outline(sheet: FactSheet, items: list[dict], vehicle: str) -> dict:
    """Group items into themed sections + pick representative timestamps (LLM, with fallback)."""
    if settings.openai_api_key:
        try:
            return _outline_llm(items, vehicle)
        except Exception as e:
            log.warning("report_pptx: LLM outline failed (%s) — using keyword fallback", e)
    return _outline_fallback(items)


def _outline_llm(items: list[dict], vehicle: str) -> dict:
    from openai import OpenAI

    compact = [
        {
            "i": i,
            "text": it["text"],
            "timestamps": sorted({round(float(e.get("t_start", 0)), 1) for e in it["evidence"]}),
        }
        for i, it in enumerate(items)
    ]
    prompt = (
        "You organize extracted vehicle facts into a concise management report outline. "
        "Group the items into coherent THEMES (e.g. steering wheel & controls, connectivity & "
        "ports, display/infotainment UI, climate & comfort, seats & interior, driver assistance "
        "& drive modes, instruments, user-journey interactions). Order themes logically. For each "
        "theme pick 1-2 key_timestamps (from the items' own timestamps) whose frames best "
        "illustrate the most items together — minimize images. Write a one-sentence summary per "
        "theme. Use the report language matching the facts.\n\n"
        "Return ONLY JSON: {\"title\": \"...\", \"sections\": [{\"title\": \"...\", "
        "\"summary\": \"...\", \"item_indices\": [int], \"key_timestamps\": [number]}]}\n\n"
        f"Vehicle: {vehicle or 'unknown'}\nItems:\n"
        + json.dumps(compact, ensure_ascii=False)
    )
    client = OpenAI(api_key=settings.openai_api_key)
    resp = client.chat.completions.create(
        model=settings.openai_model,
        response_format={"type": "json_object"},
        messages=[{"role": "user", "content": prompt}],
    )
    data = json.loads(resp.choices[0].message.content or "{}")
    if not data.get("sections"):
        raise ValueError("empty outline")
    return data


def _outline_fallback(items: list[dict]) -> dict:
    buckets: dict[str, list[int]] = {title: [] for title, _ in _FALLBACK_THEMES}
    buckets["Sonstiges"] = []
    for i, it in enumerate(items):
        low = it["text"].lower()
        placed = False
        for title, kws in _FALLBACK_THEMES:
            if any(k in low for k in kws):
                buckets[title].append(i)
                placed = True
                break
        if not placed:
            buckets["Sonstiges"].append(i)
    sections = []
    for title, idxs in buckets.items():
        if not idxs:
            continue
        ts = sorted({round(float(e.get("t_start", 0)), 1)
                     for i in idxs for e in items[i]["evidence"]})
        sections.append({
            "title": title,
            "summary": "",
            "item_indices": idxs,
            "key_timestamps": _dedupe_ts(ts)[:2],
        })
    return {"title": "Fahrzeug-Analyse", "sections": sections}


def _dedupe_ts(timestamps: list[float], bucket: float = 2.0) -> list[float]:
    """Collapse near-identical timestamps so we reuse one frame for nearby evidence."""
    out: list[float] = []
    for t in sorted(timestamps):
        if not out or abs(t - out[-1]) >= bucket:
            out.append(t)
    return out


def _frame_loader(video_path: Path | None):
    """Return a cached t -> jpeg-bytes loader (None when no video / no OpenCV)."""
    cache: dict[int, bytes | None] = {}
    cap = None
    if video_path and video_path.exists():
        try:
            import cv2

            cap = cv2.VideoCapture(str(video_path))
        except Exception as e:
            log.warning("report_pptx: OpenCV unavailable (%s)", e)

    def load(t: float) -> bytes | None:
        if cap is None:
            return None
        key = int(round(t * 10))
        if key in cache:
            return cache[key]
        import cv2

        cap.set(cv2.CAP_PROP_POS_MSEC, t * 1000.0)
        ok, frame = cap.read()
        jpeg = None
        if ok and frame is not None:
            ok_enc, buf = cv2.imencode(".jpg", frame)
            if ok_enc:
                jpeg = buf.tobytes()
        cache[key] = jpeg
        return jpeg

    return load, (cap.release if cap is not None else (lambda: None))


# ---------------------------------------------------------------------------
# Slide building
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ITEMS_PER_SLIDE = 6


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _title_bar(slide, text: str):
    tf = _textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = DARK


def _add_title_slide(prs, title: str, vehicle: str, job: Job):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = _textbox(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
    if vehicle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = vehicle
        r2.font.size = Pt(24); r2.font.color.rgb = ACCENT
    sub = _textbox(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
    ps = sub.paragraphs[0]
    rs = ps.add_run()
    rs.text = f"Quelle: {job.filename} · erstellt aus Videoanalyse"
    rs.font.size = Pt(11); rs.font.color.rgb = GRAY


def _add_text_slide(prs, title: str, body: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, title)
    tf = _textbox(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = body
    r.font.size = Pt(16); r.font.color.rgb = DARK


def _add_agenda_slide(prs, sections: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Inhalt")
    tf = _textbox(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4))
    for i, s in enumerate(sections):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = f"{i + 1}.  {s['title']}"
        r.font.size = Pt(18); r.font.color.rgb = DARK
        p.space_after = Pt(6)


def _evidence_sources(ev: list[dict]) -> str:
    """Full source string for the Quellenverzeichnis (all evidence of one item)."""
    if not ev:
        return "(kein Beleg)"
    parts = []
    for e in ev:
        ts = _fmt_ts(float(e.get("t_start", 0)))
        end = e.get("t_end")
        if end:
            ts += f"–{_fmt_ts(float(end))}"
        quote = (e.get("quote") or "").strip()
        note = (e.get("note") or "").strip()
        extra = f"„{quote}“" if quote else note
        parts.append(ts + (f" · {extra}" if extra else ""))
    return "; ".join(parts)


def _add_section_slides(prs, section: dict, items: list[dict], load_frame):
    idxs = [i for i in section.get("item_indices", []) if 0 <= i < len(items)]
    if not idxs:
        return
    key_ts = _dedupe_ts([float(t) for t in section.get("key_timestamps", [])])[:2]
    # Fallback image: first evidence timestamp in the section.
    if not key_ts:
        for i in idxs:
            if items[i]["evidence"]:
                key_ts = [float(items[i]["evidence"][0].get("t_start", 0))]
                break

    pages = [idxs[i:i + ITEMS_PER_SLIDE] for i in range(0, len(idxs), ITEMS_PER_SLIDE)]
    for pi, page in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        suffix = "" if len(pages) == 1 else f" ({pi + 1}/{len(pages)})"
        _title_bar(slide, section["title"] + suffix)

        if section.get("summary") and pi == 0:
            cap = _textbox(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5))
            rp = cap.paragraphs[0].add_run()
            rp.text = section["summary"]; rp.font.size = Pt(13); rp.font.italic = True
            rp.font.color.rgb = GRAY

        has_img = pi == 0 and bool(key_ts)
        body_w = Inches(7.0) if has_img else Inches(12.1)
        tf = _textbox(slide, Inches(0.6), Inches(1.8), body_w, Inches(5.2))
        for j, i in enumerate(page):
            it = items[i]
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = "• " + it["text"] + " "
            r.font.size = Pt(15); r.font.color.rgb = DARK
            # Small, unobtrusive reference number into the Quellenverzeichnis.
            ref = p.add_run(); ref.text = f"[{i + 1}]"
            ref.font.size = Pt(9); ref.font.color.rgb = GRAY
            p.space_after = Pt(6)

        if has_img:
            top = Inches(1.8)
            for t in key_ts:
                jpeg = load_frame(t)
                if not jpeg:
                    continue
                slide.shapes.add_picture(BytesIO(jpeg), Inches(7.9), top, width=Inches(4.9))
                lab = _textbox(slide, Inches(7.9), top + Inches(2.55), Inches(4.9), Inches(0.3))
                rl = lab.paragraphs[0].add_run()
                rl.text = f"⟶ {_fmt_ts(t)}"; rl.font.size = Pt(9); rl.font.color.rgb = GRAY
                top = top + Inches(2.9)


def _embed_movie(slide, video_path: Path, load_frame, left, top, width, height) -> bool:
    """Embed the source video into a slide (with a poster frame). Best-effort."""
    try:
        poster = load_frame(1.0) or load_frame(0.0)
        mime = "video/quicktime" if str(video_path).lower().endswith(".mov") else "video/mp4"
        slide.shapes.add_movie(
            str(video_path), left, top, width, height,
            poster_frame_image=BytesIO(poster) if poster else None,
            mime_type=mime,
        )
        return True
    except Exception as e:
        log.warning("report_pptx: could not embed video (%s)", e)
        return False


def _fill_entries(tf, chunk: list[str]):
    for k, line in enumerate(chunk):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(11); r.font.color.rgb = DARK
        p.space_after = Pt(4)


def _add_sources_slides(prs, items: list[dict], job: Job, video_path: Path | None, load_frame):
    """Quellenverzeichnis: numbered sources, with the source video embedded under the heading."""
    entries = [
        f"[{i + 1}] {it['text'][:70]} — {_evidence_sources(it['evidence'])}"
        for i, it in enumerate(items)
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Quellenverzeichnis")

    has_video = False
    if video_path:
        has_video = _embed_movie(
            slide, video_path, load_frame, Inches(0.6), Inches(1.4), Inches(6.4), Inches(3.6)
        )
        if has_video:
            cap = _textbox(slide, Inches(0.6), Inches(5.05), Inches(6.4), Inches(0.4))
            rc = cap.paragraphs[0].add_run()
            rc.text = f"Video: {job.filename}"; rc.font.size = Pt(10); rc.font.color.rgb = GRAY

    if has_video:
        tf = _textbox(slide, Inches(7.2), Inches(1.4), Inches(5.5), Inches(5.6))
        _fill_entries(tf, entries[:9])
        rest = entries[9:]
    else:
        tf = _textbox(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4))
        _fill_entries(tf, entries[:16])
        rest = entries[16:]

    for c in range(0, len(rest), 16):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _title_bar(s, "Quellenverzeichnis (Forts.)")
        tf2 = _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4))
        _fill_entries(tf2, rest[c:c + 16])

    # Short methodology footer slide.
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(s, "Methodik")
    tf3 = _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4))
    lines = [
        f"Quelle: Video „{job.filename}“ — siehe eingebettetes Video im Quellenverzeichnis.",
        "Belege je Fakt als Zeitstempel (Sprungmarke ins Video) bzw. wörtliches Zitat.",
        "Erstellt per KI-Frame-Sampling und Voiceover-Transkript.",
        "Hinweis: KI-generiert — vor Veröffentlichung stichprobenartig prüfen.",
    ]
    for i, ln in enumerate(lines):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        r = p.add_run(); r.text = ln; r.font.size = Pt(12); r.font.color.rgb = GRAY
        p.space_after = Pt(8)


def generate_pptx(job: Job, sheet: FactSheet, video_path: Path | None) -> bytes:
    """Build a management-ready PowerPoint deck from a fact-sheet."""
    items = _items_from_sheet(sheet)
    vehicle = _vehicle_label(job, sheet)
    outline = _outline(sheet, items, vehicle)
    sections = outline.get("sections", [])
    load_frame, release = _frame_loader(video_path)
    try:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        title = outline.get("title") or "Fahrzeug-Analyse"
        _add_title_slide(prs, title, vehicle, job)
        if sheet.summary.strip():
            _add_text_slide(prs, "Zusammenfassung", sheet.summary.strip())
        if sections:
            _add_agenda_slide(prs, sections)
        for s in sections:
            _add_section_slides(prs, s, items, load_frame)
        _add_sources_slides(prs, items, job, video_path, load_frame)

        out = BytesIO()
        prs.save(out)
        return out.getvalue()
    finally:
        release()
